# -*- coding: utf-8 -*-
"""
build_slides.py
---------------
Tao toan bo 11 slides PPTX cho RBL-3 Proposal Defense tu dau.
Design: dark theme (#0d1117), cyan (#00bcd4), amber (#ffc107), green (#00e676)

Yeu cau: pip install python-pptx
Chay:    python build_slides.py
"""

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── OUTPUT ─────────────────────────────────────────────────────
OUT = pathlib.Path(
    r"D:\pal\RBL-SWT-main_3\SWT301-SE1944-Group6\presentation"
    r"\GPT4o_RBL3_Defense_Final_v2.pptx"
)

# ── COLORS ─────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x11, 0x17)   # dark navy
C_BG2       = RGBColor(0x16, 0x1B, 0x27)   # card bg
C_CYAN      = RGBColor(0x00, 0xBC, 0xD4)   # cyan accent
C_AMBER     = RGBColor(0xFF, 0xC1, 0x07)   # amber highlight
C_GREEN     = RGBColor(0x00, 0xE6, 0x76)   # green positive
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # white text
C_GRAY      = RGBColor(0x8B, 0x94, 0x9E)   # muted text
C_LIGHTGRAY = RGBColor(0xC9, 0xD1, 0xD9)  # light text
C_RED       = RGBColor(0xFF, 0x5C, 0x5C)   # warning
C_PURPLE    = RGBColor(0xBB, 0x86, 0xFC)   # purple

# ── SLIDE SIZE: 16:9 ───────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Blank layout


# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def new_slide():
    sl = prs.slides.add_slide(blank_layout)
    return sl


def set_bg(slide, color: RGBColor):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor = None,
             line_color: RGBColor = None, line_width_pt=0.5, radius=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def para(tf, text, *, size=18, bold=False, italic=False,
         color: RGBColor = None, align=PP_ALIGN.LEFT,
         space_before=0, space_after=0, clear=False, level=0):
    """Add or get a paragraph in a text frame."""
    if clear:
        tf.clear()
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color or C_WHITE
        run.font.name = "Calibri"
    return p


def add_divider(slide, x, y, w, color=C_CYAN, height=0.03):
    """Add a horizontal divider line."""
    rect = add_rect(slide, x, y, w, height, fill_color=color)
    return rect


def add_slide_header(slide, title, subtitle=None,
                     speaker=None, time_str=None,
                     title_color=C_CYAN):
    """Standard header: colored bar + title."""
    # Top accent bar
    add_rect(slide, 0, 0, 13.33, 0.06, fill_color=title_color)

    # Title area background
    add_rect(slide, 0, 0.06, 13.33, 1.1, fill_color=C_BG2)

    # Title text
    _, tf = add_textbox(slide, 0.35, 0.1, 10.0, 0.9)
    para(tf, title, size=32, bold=True, color=title_color,
         align=PP_ALIGN.LEFT, clear=True)

    # Subtitle
    if subtitle:
        _, tf2 = add_textbox(slide, 0.35, 0.78, 9.0, 0.38)
        para(tf2, subtitle, size=14, color=C_GRAY, clear=True)

    # Speaker badge (top right)
    if speaker:
        badge_txt = f"  {speaker}  "
        if time_str:
            badge_txt += f"| {time_str}"
        add_rect(slide, 10.8, 0.22, 2.2, 0.42, fill_color=title_color)
        _, btf = add_textbox(slide, 10.82, 0.24, 2.18, 0.38)
        para(btf, badge_txt, size=13, bold=True,
             color=C_BG, align=PP_ALIGN.CENTER, clear=True)


def add_card(slide, x, y, w, h, title=None, title_color=C_CYAN,
             bg=C_BG2, border_color=None):
    """Add a content card with optional title bar."""
    add_rect(slide, x, y, w, h, fill_color=bg,
             line_color=border_color or title_color, line_width_pt=0.8)
    if title:
        add_rect(slide, x, y, w, 0.35, fill_color=title_color)
        _, tf = add_textbox(slide, x + 0.12, y + 0.03, w - 0.15, 0.3)
        para(tf, title, size=13, bold=True,
             color=C_BG, align=PP_ALIGN.LEFT, clear=True)
    return y + (0.35 if title else 0)


def slide_num_label(slide, num, total=11):
    """Small slide number at bottom right."""
    _, tf = add_textbox(slide, 12.4, 7.22, 0.8, 0.22)
    para(tf, f"{num}/{total}", size=11, color=C_GRAY,
         align=PP_ALIGN.RIGHT, clear=True)

    # Bottom accent line
    add_rect(slide, 0, 7.38, 13.33, 0.06, fill_color=C_BG2)
    add_rect(slide, 0, 7.42, 13.33, 0.08, fill_color=C_CYAN)


def bullet(tf, text, *, size=16, color=C_LIGHTGRAY, indent=0.0,
           bold=False, italic=False, marker="•", space_before=4):
    """Add a bullet paragraph."""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(space_before)
    # Marker run
    if marker:
        r1 = p.add_run()
        r1.text = marker + "  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = C_CYAN
        r1.font.name = "Calibri"
        r1.font.bold = True
    # Content run
    r2 = p.add_run()
    r2.text = text
    r2.font.size = Pt(size)
    r2.font.color.rgb = color
    r2.font.name = "Calibri"
    r2.font.bold = bold
    r2.font.italic = italic
    return p


def kv_line(tf, key, val, *, key_size=14, val_size=14,
            key_color=C_CYAN, val_color=C_WHITE, space=3):
    """Key: Value line in a text frame."""
    p = tf.add_paragraph()
    p.space_before = Pt(space)
    r1 = p.add_run()
    r1.text = key + ":  "
    r1.font.size = Pt(key_size)
    r1.font.bold = True
    r1.font.color.rgb = key_color
    r1.font.name = "Calibri"
    r2 = p.add_run()
    r2.text = val
    r2.font.size = Pt(val_size)
    r2.font.color.rgb = val_color
    r2.font.name = "Calibri"
    return p


# ═══════════════════════════════════════════════════════════════
# S1 — TITLE & SCOPE
# ═══════════════════════════════════════════════════════════════
def build_s1():
    sl = new_slide()
    set_bg(sl, C_BG)

    # Top gradient bar
    add_rect(sl, 0, 0, 13.33, 0.12, fill_color=C_CYAN)

    # Center hero area
    add_rect(sl, 0.5, 0.6, 12.33, 3.8, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=1.5)
    add_rect(sl, 0.5, 0.6, 0.12, 3.8, fill_color=C_CYAN)

    # Main title
    _, tf = add_textbox(sl, 0.9, 0.75, 11.5, 1.8)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Danh gia First-Attempt Zero-Shot"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "GPT-4o Unit Test Generation"
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.color.rgb = C_CYAN
    r2.font.name = "Calibri"

    p3 = tf.add_paragraph()
    p3.space_before = Pt(6)
    r3 = p3.add_run()
    r3.text = "tren Java/Python Functions voi CC 5-15"
    r3.font.size = Pt(24)
    r3.font.bold = False
    r3.font.color.rgb = C_LIGHTGRAY
    r3.font.name = "Calibri"

    # Subtitle line
    add_divider(sl, 0.9, 2.75, 6.0, color=C_CYAN, height=0.04)

    _, tf2 = add_textbox(sl, 0.9, 2.88, 11.0, 0.4)
    para(tf2, "Controlled Empirical Study  |  SWT301 — Software Testing  |  2026",
         size=15, color=C_GRAY, clear=True)

    # Scope cards row
    cards = [
        ("50", "Units\n25 Java + 25 Python", C_CYAN),
        ("CC 5-15", "Cyclomatic\nComplexity Range", C_AMBER),
        ("Zero-Shot", "First Attempt\nNo Repair", C_GREEN),
        ("3 RQs", "Coverage • Mutation\n• Correlation", C_PURPLE),
    ]
    cx = 0.9
    for val, lbl, col in cards:
        add_rect(sl, cx, 3.7, 2.8, 1.5, fill_color=C_BG,
                 line_color=col, line_width_pt=1.5)
        add_rect(sl, cx, 3.7, 2.8, 0.06, fill_color=col)
        _, tf3 = add_textbox(sl, cx + 0.1, 3.8, 2.6, 0.7)
        para(tf3, val, size=28, bold=True, color=col,
             align=PP_ALIGN.CENTER, clear=True)
        _, tf4 = add_textbox(sl, cx + 0.1, 4.52, 2.6, 0.6)
        para(tf4, lbl, size=11, color=C_GRAY, align=PP_ALIGN.CENTER, clear=True)
        cx += 3.0

    # Team members bottom
    _, tf5 = add_textbox(sl, 0.5, 5.45, 12.33, 0.35)
    para(tf5, "Nhom: Nghi (PL)  |  Vinh (DG)  |  Thai (LR)  |  Quan (MS)  |  Van (RW)",
         size=13, color=C_GRAY, align=PP_ALIGN.CENTER, clear=True)

    # bottom bar
    add_rect(sl, 0, 7.38, 13.33, 0.12, fill_color=C_CYAN)
    _, tbn = add_textbox(sl, 12.0, 7.25, 1.2, 0.25)
    para(tbn, "1 / 11", size=11, color=C_GRAY, align=PP_ALIGN.RIGHT, clear=True)


# ═══════════════════════════════════════════════════════════════
# S2 — PROBLEM → GAP
# ═══════════════════════════════════════════════════════════════
def build_s2():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S2  Problem  GAP", "Confounders trong literature -> Bounded GAP",
                     speaker="Vinh", time_str="35s")
    slide_num_label(sl, 2)

    # Three boxes: Prior Work | Confounders | GAP
    boxes = [
        ("PRIOR WORK", C_CYAN, [
            "TestPilot: median branch",
            "  coverage 52,8%",
            "  (Schafer 2024, IEEE TSE)",
            "",
            "AgoneTest: highest branch",
            "  coverage 79,8%",
            "  (Lops 2025, ASE)",
            "",
            "Compilation AgoneTest: 36%",
        ]),
        ("CONFOUNDERS", C_AMBER, [
            "Dataset khac nhau",
            "(JS/Python vs Java class)",
            "",
            "Model khac nhau",
            "(gpt-3.5 vs gpt-4o-mini)",
            "",
            "Prompt khac nhau",
            "(zero-shot vs few-shot)",
            "",
            "Metric: median vs highest",
        ]),
        ("BOUNDED GAP", C_GREEN, [
            "Trong 10 papers ra soat:",
            "",
            "Chua co nghien cuu danh",
            "gia dong thoi mau can bang",
            "Java-Python loc truoc theo",
            "CC 5-15",
            "",
            "-> Can controlled protocol",
        ]),
    ]

    bx = 0.35
    for title, col, lines in boxes:
        add_rect(sl, bx, 1.35, 4.1, 5.7, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)
        add_rect(sl, bx, 1.35, 4.1, 0.42, fill_color=col)

        _, th = add_textbox(sl, bx + 0.15, 1.38, 3.8, 0.35)
        para(th, title, size=15, bold=True, color=C_BG, align=PP_ALIGN.CENTER, clear=True)

        _, tf = add_textbox(sl, bx + 0.18, 1.88, 3.75, 4.8)
        tf.word_wrap = True
        first = True
        for line in lines:
            if first:
                para(tf, line, size=14, color=C_WHITE if col != C_GREEN else C_WHITE,
                     clear=True)
                first = False
            else:
                para(tf, line, size=14, color=C_LIGHTGRAY)
        bx += 4.33

    # Arrow labels between boxes
    _, ta1 = add_textbox(sl, 4.3, 4.1, 0.5, 0.5)
    para(ta1, "->", size=22, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER, clear=True)
    _, ta2 = add_textbox(sl, 8.63, 4.1, 0.5, 0.5)
    para(ta2, "->", size=22, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, clear=True)

    # Bottom callout
    add_rect(sl, 0.35, 7.05, 12.63, 0.28, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=1.0)
    _, tbot = add_textbox(sl, 0.5, 7.07, 12.3, 0.24)
    para(tbot, "GAP-D: Chua co nghien cuu nao trong 10 papers dung mau can bang Java/Python loc truoc CC 5-15",
         size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# S3 — KEY EVIDENCE TABLE
# ═══════════════════════════════════════════════════════════════
def build_s3():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S3  Key Evidence", "5 paper lien quan truc tiep nhat tu SLR 10 papers",
                     speaker="Vinh", time_str="45s")
    slide_num_label(sl, 3)

    # Table header
    cols = [2.8, 2.2, 2.2, 1.8, 4.0]
    headers = ["Paper", "Branch Cov.", "Mutation", "Compilability", "Ghi chu quan trong"]
    hx = 0.35
    hy = 1.45
    for i, (w, h) in enumerate(zip(cols, headers)):
        add_rect(sl, hx, hy, w - 0.05, 0.42, fill_color=C_CYAN)
        _, tf = add_textbox(sl, hx + 0.08, hy + 0.04, w - 0.1, 0.34)
        para(tf, h, size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER, clear=True)
        hx += w

    # Table rows
    rows = [
        ("TestPilot\n(Schafer 2024, TSE)", "Median\n52,8%", "NR", "NR",
         "gpt-3.5-turbo; JS/Python\nfunction-level; zero-shot",
         C_CYAN, False),
        ("AgoneTest\n(Lops 2025, ASE)", "Highest\n79,8%", "Highest\n89,2%", "36%",
         "Max nhieu model/prompt\nKHONG phai GPT-4o",
         C_AMBER, True),
        ("GEM\n(Celik 2026, CIbSE)", "84,30%->88,81%", "72,71%->81,85%", "NR",
         "Ket qua SAU repair\n& strengthening",
         C_GRAY, False),
        ("ASTER\n(Pan 2024, ICSE)", "77,2%", "NR", "NR",
         "Static analysis +\nerror remediation",
         C_GRAY, False),
        ("ChatTester\n(Yuan 2024, FSE)", ">82%", ">65%", "Tang 34,3%",
         "Java class-level;\nmulti-turn repair",
         C_GRAY, False),
    ]

    ry = hy + 0.44
    for rdata in rows:
        cells = rdata[:5]
        row_col = rdata[5]
        is_warn = rdata[6]
        rx = 0.35
        row_h = 0.93
        bg = RGBColor(0x1C, 0x24, 0x32) if not is_warn else RGBColor(0x2A, 0x22, 0x0A)
        for i, (w, cell) in enumerate(zip(cols, cells)):
            add_rect(sl, rx, ry, w - 0.05, row_h, fill_color=bg,
                     line_color=RGBColor(0x30, 0x40, 0x50), line_width_pt=0.5)
            _, tf = add_textbox(sl, rx + 0.08, ry + 0.05, w - 0.12, row_h - 0.08)
            tf.word_wrap = True
            col_color = row_col if i in (1, 2) and row_col != C_GRAY else C_LIGHTGRAY
            if i == 0:
                col_color = C_WHITE
            para(tf, cell, size=12, color=col_color,
                 align=PP_ALIGN.CENTER, bold=(i == 0), clear=True)
            rx += w
        ry += row_h + 0.03

    # Bottom note
    add_rect(sl, 0.35, 6.85, 12.63, 0.3, fill_color=RGBColor(0x1A, 0x1A, 0x1A),
             line_color=C_AMBER, line_width_pt=0.8)
    _, tbot = add_textbox(sl, 0.5, 6.88, 12.3, 0.26)
    para(tbot, "Chi dung cac so nay lam BOI CANH — khong gop meta-analysis, khong suy ra universal threshold",
         size=12, bold=True, color=C_AMBER, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# S4 — RQ & STATISTICAL TESTS
# ═══════════════════════════════════════════════════════════════
def build_s4():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S4  Research Questions & Statistical Tests",
                     "Slide quan trong nhat — GV se hoi nhieu",
                     speaker="Quan", time_str="50s", title_color=C_AMBER)
    slide_num_label(sl, 4)

    rqs = [
        ("RQ1", "Branch Coverage vs Practical Target", C_CYAN,
         "H0: median(BC) <= 75%  |  H1: median(BC) > 75%",
         "Test: Exact Sign Test, one-sided, right-tailed",
         "Multiplicity: Holm correction, 2 tests (Java + Python)",
         "Threshold: 75%  =  Case 3 — Pre-registered Practical Target"),
        ("RQ2", "GPT-4o vs Randoop 4.3.4 (Java Mutation Score)", C_AMBER,
         "d_i = mutation_GPT_i - mutation_Randoop_i",
         "H0: median(d) <= 0  |  H1: median(d) > 0",
         "Test: Paired Wilcoxon Signed-Rank, one-sided, right-tailed",
         "Bao cao: effect size + 95% CI du khi GPT-4o thap hon"),
        ("RQ3", "CC vs Test Quality (Correlation)", C_GREEN,
         "4 cap: Java/Python x Branch Coverage/Mutation Score",
         "H0: rho >= 0  |  H1: rho < 0",
         "Test: Spearman, one-sided, left-tailed",
         "Holm correction, 4 tests  |  EXPLORATORY — khong suy nhan qua"),
    ]

    ry = 1.38
    for rq_id, rq_title, col, l1, l2, l3, l4 in rqs:
        add_rect(sl, 0.35, ry, 12.63, 1.73, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)
        add_rect(sl, 0.35, ry, 1.1, 1.73, fill_color=col)

        _, tid = add_textbox(sl, 0.37, ry + 0.55, 1.08, 0.6)
        para(tid, rq_id, size=22, bold=True, color=C_BG,
             align=PP_ALIGN.CENTER, clear=True)

        _, tt = add_textbox(sl, 1.55, ry + 0.08, 11.2, 0.38)
        para(tt, rq_title, size=15, bold=True, color=col, clear=True)

        _, tb = add_textbox(sl, 1.55, ry + 0.48, 11.2, 1.2)
        tb.word_wrap = True
        para(tb, l1, size=13, color=C_WHITE, clear=True)
        para(tb, l2, size=12, color=C_LIGHTGRAY)
        para(tb, l3, size=12, color=C_LIGHTGRAY)
        r = tb.add_paragraph()
        r.space_before = Pt(3)
        ru = r.add_run()
        ru.text = l4
        ru.font.size = Pt(12)
        ru.font.bold = True
        ru.font.color.rgb = col
        ru.font.name = "Calibri"

        ry += 1.83

    # Family alpha note
    add_rect(sl, 0.35, 7.0, 12.63, 0.3, fill_color=C_BG2,
             line_color=C_PURPLE, line_width_pt=1.0)
    _, tfa = add_textbox(sl, 0.5, 7.03, 12.3, 0.26)
    para(tfa, "Family-wise alpha = 0.05  |  Chi dung 'reject H0' hoac 'fail to reject H0'  |  Khong thay threshold sau khi co data",
         size=12, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# S5 — PIPELINE & MODEL SPECIFICATION
# ═══════════════════════════════════════════════════════════════
def build_s5():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S5  Pipeline & Model Specification",
                     "Fixed before run — khong thay doi sau khi co data",
                     speaker="Thai", time_str="50s", title_color=C_GREEN)
    slide_num_label(sl, 5)

    # Pipeline flow: 6 steps
    steps = [
        ("1\nFreeze\nManifest", C_CYAN),
        ("2\nLizard CC\nFilter", C_CYAN),
        ("3\nRender\nPrompt", C_AMBER),
        ("4\nGPT-4o\nAPI Call", C_GREEN),
        ("5\nCompile &\nMeasure", C_PURPLE),
        ("6\nStats\nAnalysis", C_RED),
    ]
    sx = 0.3
    for i, (lbl, col) in enumerate(steps):
        add_rect(sl, sx, 1.42, 1.95, 1.05, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)
        add_rect(sl, sx, 1.42, 1.95, 0.08, fill_color=col)
        _, tf = add_textbox(sl, sx + 0.08, 1.5, 1.8, 0.95)
        para(tf, lbl, size=13, bold=True, color=col,
             align=PP_ALIGN.CENTER, clear=True)
        if i < 5:
            _, ta = add_textbox(sl, sx + 1.97, 1.82, 0.32, 0.32)
            para(ta, "->", size=16, bold=True, color=C_GRAY,
                 align=PP_ALIGN.CENTER, clear=True)
        sx += 2.18

    # Model spec card (left)
    add_rect(sl, 0.35, 2.72, 5.8, 4.3, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=1.2)
    add_rect(sl, 0.35, 2.72, 5.8, 0.4, fill_color=C_GREEN)
    _, tth = add_textbox(sl, 0.5, 2.75, 5.5, 0.34)
    para(tth, "MODEL SPECIFICATION", size=14, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)

    _, tspec = add_textbox(sl, 0.5, 3.22, 5.5, 3.7)
    tspec.word_wrap = True
    specs = [
        ("Model", "gpt-4o-2024-08-06"),
        ("Endpoint", "POST /v1/responses"),
        ("SDK", "OpenAI Python 2.43.0"),
        ("Temperature", "0.2"),
        ("top_p", "1.0"),
        ("max_output_tokens", "2,048"),
        ("Responses/unit", "1  (NO regeneration)"),
        ("Retry", "Infrastructure only, before response"),
    ]
    first = True
    for k, v in specs:
        if first:
            kv_line(tspec, k, v, key_size=13, val_size=13, space=0)
            first = False
        else:
            kv_line(tspec, k, v, key_size=13, val_size=13, space=4)

    # Toolchain card (right)
    add_rect(sl, 6.5, 2.72, 6.48, 4.3, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=1.2)
    add_rect(sl, 6.5, 2.72, 6.48, 0.4, fill_color=C_CYAN)
    _, tth2 = add_textbox(sl, 6.65, 2.75, 6.1, 0.34)
    para(tth2, "TOOLCHAIN", size=14, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)

    _, ttool = add_textbox(sl, 6.65, 3.22, 6.1, 3.7)
    ttool.word_wrap = True
    tools = [
        ("Lizard", "1.23.0  — CC measurement"),
        ("JaCoCo", "0.8.13  — Java branch coverage"),
        ("coverage.py", "7.14.2  — Python branch coverage"),
        ("PIT", "1.19.1  — Java mutation score"),
        ("pytest-mutagen", "1.3  — Python mutation score"),
        ("Randoop", "4.3.4  — Java baseline"),
        ("JUnit", "5.11.4  |  JDK 17"),
        ("SciPy", "1.18.0  |  pandas 3.0.3"),
    ]
    first = True
    for k, v in tools:
        if first:
            kv_line(ttool, k, v, key_size=13, val_size=13, space=0)
            first = False
        else:
            kv_line(ttool, k, v, key_size=13, val_size=13, space=4)


# ═══════════════════════════════════════════════════════════════
# S6 — DATASET & BASELINE
# ═══════════════════════════════════════════════════════════════
def build_s6():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S6  Dataset & Baseline",
                     "50 units, CC 5-15, provenance gate",
                     speaker="Vinh", time_str="40s")
    slide_num_label(sl, 6)

    # Dataset tree visual
    # Root node
    add_rect(sl, 4.85, 1.42, 3.63, 0.55, fill_color=C_CYAN)
    _, tr = add_textbox(sl, 4.9, 1.46, 3.53, 0.45)
    para(tr, "DATASET  N = 50", size=17, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)

    # Branch lines (drawn as thin rects)
    add_rect(sl, 6.65, 1.97, 0.04, 0.4, fill_color=C_GRAY)   # down
    add_rect(sl, 2.65, 2.37, 4.05, 0.04, fill_color=C_GRAY)   # horizontal
    add_rect(sl, 6.65, 2.37, 4.05, 0.04, fill_color=C_GRAY)
    add_rect(sl, 2.65, 2.37, 0.04, 0.4, fill_color=C_GRAY)    # left branch
    add_rect(sl, 10.66, 2.37, 0.04, 0.4, fill_color=C_GRAY)   # right branch

    # Java branch
    add_rect(sl, 1.2, 2.77, 3.0, 0.55, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=2.0)
    _, tj = add_textbox(sl, 1.25, 2.79, 2.9, 0.48)
    para(tj, "JAVA  25 units", size=16, bold=True,
         color=C_AMBER, align=PP_ALIGN.CENTER, clear=True)

    add_rect(sl, 1.2, 3.38, 3.0, 0.45, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.0)
    _, tr1 = add_textbox(sl, 1.3, 3.4, 2.8, 0.38)
    para(tr1, "Randoop 4.3.4 baseline\n60s/unit | seed=20260621",
         size=12, color=C_LIGHTGRAY, align=PP_ALIGN.CENTER, clear=True)

    add_rect(sl, 1.2, 3.89, 3.0, 0.45, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.0)
    _, tr2 = add_textbox(sl, 1.3, 3.91, 2.8, 0.38)
    para(tr2, "Benchmark: CLASSES2TEST\n(Lops 2025, ASE CORE A*)",
         size=12, color=C_LIGHTGRAY, align=PP_ALIGN.CENTER, clear=True)

    # Python branch
    add_rect(sl, 9.13, 2.77, 3.0, 0.55, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=2.0)
    _, tp = add_textbox(sl, 9.18, 2.79, 2.9, 0.48)
    para(tp, "PYTHON  25 units", size=16, bold=True,
         color=C_GREEN, align=PP_ALIGN.CENTER, clear=True)

    add_rect(sl, 9.13, 3.38, 3.0, 0.45, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=1.0)
    _, tp1 = add_textbox(sl, 9.23, 3.4, 2.8, 0.38)
    para(tp1, "No traditional baseline",
         size=12, color=C_LIGHTGRAY, align=PP_ALIGN.CENTER, clear=True)

    add_rect(sl, 9.13, 3.89, 3.0, 0.45, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=1.0)
    _, tp2 = add_textbox(sl, 9.23, 3.91, 2.8, 0.38)
    para(tp2, "Method: TestPilot + ASTER\n(IEEE TSE + ICSE CORE A*)",
         size=12, color=C_LIGHTGRAY, align=PP_ALIGN.CENTER, clear=True)

    # Shared CC gate at bottom
    add_rect(sl, 4.2, 4.55, 4.95, 0.55, fill_color=C_BG2,
             line_color=C_PURPLE, line_width_pt=2.0)
    add_rect(sl, 4.2, 4.55, 4.95, 0.08, fill_color=C_PURPLE)
    _, tcc = add_textbox(sl, 4.3, 4.64, 4.75, 0.4)
    para(tcc, "5 <= CC <= 15  |  Lizard 1.23.0",
         size=15, bold=True, color=C_WHITE,
         align=PP_ALIGN.CENTER, clear=True)

    # Provenance gate
    add_rect(sl, 4.2, 5.22, 4.95, 0.55, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=1.5)
    _, tpg = add_textbox(sl, 4.3, 5.24, 4.75, 0.5)
    para(tpg, "Provenance Gate:\nsource • license • commit SHA • hash",
         size=13, color=C_CYAN, align=PP_ALIGN.CENTER, clear=True)

    # Inclusion rules card
    add_rect(sl, 0.35, 5.9, 12.63, 1.2, fill_color=C_BG2,
             line_color=C_GRAY, line_width_pt=0.8)
    _, tic = add_textbox(sl, 0.5, 5.96, 12.3, 0.26)
    para(tic, "INCLUSION RULES:", size=13, bold=True, color=C_CYAN, clear=True)
    _, tic2 = add_textbox(sl, 0.5, 6.22, 12.3, 0.8)
    para(tic2,
         "Public source • has license • build/import standalone • no network/DB/clock deps"
         " • 2 reviewers • sort by SHA-256(unit_id) • top 25/language",
         size=12, color=C_LIGHTGRAY, clear=True)


# ═══════════════════════════════════════════════════════════════
# S7 — WORK PACKAGES
# ═══════════════════════════════════════════════════════════════
def build_s7():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S7  Work Packages",
                     "1 package — 1 owner chinh",
                     speaker="Van", time_str="40s", title_color=C_PURPLE)
    slide_num_label(sl, 7)

    owners = [
        ("VINH", "DG", "Dataset & Preflight\n50 units + manifest\n+ provenance gate", C_CYAN),
        ("THAI", "LR", "50 API Calls\nRaw logs, token cost\nPrompt hash", C_GREEN),
        ("QUAN", "MS", "25 Java + Randoop\nCoverage, Mutation\nStatistical tests", C_AMBER),
        ("VAN", "RW", "25 Python units\nFigures, validity\nobservations", C_PURPLE),
        ("NGHI", "PL", "Protocol freeze\nReproducibility gate\nFinal integration", C_RED),
    ]

    ox = 0.35
    for name, role, task, col in owners:
        add_rect(sl, ox, 1.42, 2.38, 3.0, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)
        # Top badge
        add_rect(sl, ox, 1.42, 2.38, 0.7, fill_color=col)
        _, tn = add_textbox(sl, ox + 0.05, 1.44, 2.28, 0.35)
        para(tn, name, size=20, bold=True, color=C_BG,
             align=PP_ALIGN.CENTER, clear=True)
        _, tr = add_textbox(sl, ox + 0.05, 1.78, 2.28, 0.28)
        para(tr, f"Role: {role}", size=12, bold=True, color=C_BG,
             align=PP_ALIGN.CENTER, clear=True)
        _, tt = add_textbox(sl, ox + 0.1, 2.22, 2.18, 2.1)
        tt.word_wrap = True
        para(tt, task, size=13, color=C_LIGHTGRAY,
             align=PP_ALIGN.CENTER, clear=True)
        # Arrow
        if ox < 10.0:
            _, ta = add_textbox(sl, ox + 2.4, 2.8, 0.35, 0.4)
            para(ta, "->", size=18, bold=True, color=C_GRAY,
                 align=PP_ALIGN.CENTER, clear=True)
        ox += 2.58

    # Rules
    add_rect(sl, 0.35, 4.65, 12.63, 0.65, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.0)
    _, trule = add_textbox(sl, 0.5, 4.7, 12.3, 0.6)
    trule.word_wrap = True
    para(trule, "Quy tac: LR (Thai) + MS (Quan) KHONG gop — nguoi chay experiment khong tu verify ket qua cua chinh minh.",
         size=13, bold=True, color=C_AMBER, clear=True)

    # Ownership banner
    add_rect(sl, 0.35, 5.48, 12.63, 0.55, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=1.0)
    _, tban = add_textbox(sl, 0.5, 5.52, 12.3, 0.45)
    para(tban, "Ai lam phan nao -> Tu viet phan do trong report. Khong giao 1 nguoi viet thay ca nhom.",
         size=13, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# S8 — EVALUATION PLAN
# ═══════════════════════════════════════════════════════════════
def build_s8():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S8  Evaluation Plan",
                     "Dien giai tat ca ket qua co the xay ra",
                     speaker="Quan", time_str="45s", title_color=C_AMBER)
    slide_num_label(sl, 8)

    # RQ1 outcome matrix
    _, th1 = add_textbox(sl, 0.35, 1.42, 8.0, 0.35)
    para(th1, "RQ1 — 3 kich ban ket qua:", size=15, bold=True,
         color=C_CYAN, clear=True)

    outcomes = [
        ("Ca Java VA Python\nvuot 75%", "Evidence on dinh\nca hai ngon ngu", C_GREEN),
        ("Chi MOT ngon ngu\nvuot 75%", "Ket luan:\nLanguage-specific", C_AMBER),
        ("KHONG ngon ngu nao\nvuot 75%", "Target khong duoc\nho tro  —  giu nguyen", C_RED),
    ]
    ox = 0.35
    for cond, concl, col in outcomes:
        add_rect(sl, ox, 1.88, 3.92, 1.35, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)
        add_rect(sl, ox, 1.88, 3.92, 0.06, fill_color=col)
        _, tc = add_textbox(sl, ox + 0.1, 1.96, 3.72, 0.6)
        para(tc, cond, size=13, color=C_WHITE, align=PP_ALIGN.CENTER, clear=True)
        add_divider(sl, ox + 0.15, 2.6, 3.62, color=col, height=0.03)
        _, tco = add_textbox(sl, ox + 0.1, 2.68, 3.72, 0.5)
        para(tco, concl, size=13, bold=True, color=col,
             align=PP_ALIGN.CENTER, clear=True)
        if ox < 8.0:
            _, ta = add_textbox(sl, ox + 3.94, 2.37, 0.45, 0.4)
            para(ta, "OR", size=13, bold=True, color=C_GRAY,
                 align=PP_ALIGN.CENTER, clear=True)
        ox += 4.17

    # RQ2 block
    _, th2 = add_textbox(sl, 0.35, 3.4, 8.0, 0.33)
    para(th2, "RQ2 — Bao cao paired mutation difference:", size=15,
         bold=True, color=C_AMBER, clear=True)
    add_rect(sl, 0.35, 3.8, 12.63, 0.75, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.0)
    _, trq2 = add_textbox(sl, 0.5, 3.86, 12.3, 0.65)
    trq2.word_wrap = True
    para(trq2,
         "Du GPT-4o tot hon hay khong tot hon Randoop -> bao ca hai truong hop."
         "  Bao gom: n hop le, n missing, median/IQR, paired diff,"
         " effect size, 95% CI, raw p, adjusted p.",
         size=13, color=C_LIGHTGRAY, clear=True)

    # Anti-HARKing rule
    add_rect(sl, 0.35, 4.72, 12.63, 0.6, fill_color=RGBColor(0x2A, 0x10, 0x10),
             line_color=C_RED, line_width_pt=1.5)
    _, tark = add_textbox(sl, 0.5, 4.78, 12.3, 0.5)
    para(tark, "KHONG thay threshold sau khi xem ket qua  =  HARKing  =  Vi pham tinh toan ven khoa hoc",
         size=14, bold=True, color=C_RED, align=PP_ALIGN.CENTER, clear=True)

    # Report requirements
    add_rect(sl, 0.35, 5.5, 12.63, 1.55, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=0.8)
    _, trep = add_textbox(sl, 0.5, 5.56, 12.3, 1.42)
    trep.word_wrap = True
    para(trep, "MOI bang ket qua phai co day du:", size=14,
         bold=True, color=C_CYAN, clear=True)
    para(trep, "n hop le  •  n missing  •  median / IQR  •  effect estimate  •  95% CI"
         "  •  raw p-value  •  adjusted p-value (Holm)",
         size=13, color=C_LIGHTGRAY)


# ═══════════════════════════════════════════════════════════════
# S9 — THREATS TO VALIDITY
# ═══════════════════════════════════════════════════════════════
def build_s9():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S9  Threats to Validity",
                     "4 loai bat buoc — moi threat co mitigation cu the",
                     speaker="Van", time_str="40s", title_color=C_RED)
    slide_num_label(sl, 9)

    threats = [
        ("INTERNAL", "Model Drift:\nOpenAI co the silent-update model",
         "Pin exact snapshot gpt-4o-2024-08-06\nLog version trong moi run", C_RED),
        ("EXTERNAL", "50 units, 2 languages:\nKho khai quat ra ngoai sample",
         "Claim chi trong 50 units hop le\nCC 5-15, khong overgeneralize", C_AMBER),
        ("CONSTRUCT", "Coverage != Fault Detection:\nBC cao khong co nghia test tot",
         "Co-primary: Mutation Score + Executability\nBao cao rieng, khong thay the nhau", C_PURPLE),
        ("CONCLUSION", "N nho + Range restriction:\nEstimate co the khong on dinh",
         "Bao cao effect size + 95% CI\nRQ3 la EXPLORATORY — khong suy nhan qua", C_GREEN),
    ]

    tx = 0.35
    for ttype, threat, mitigation, col in threats:
        add_rect(sl, tx, 1.42, 3.1, 5.65, fill_color=C_BG2,
                 line_color=col, line_width_pt=1.5)

        # Threat type header
        add_rect(sl, tx, 1.42, 3.1, 0.42, fill_color=col)
        _, tth = add_textbox(sl, tx + 0.08, 1.44, 2.94, 0.36)
        para(tth, ttype, size=15, bold=True, color=C_BG,
             align=PP_ALIGN.CENTER, clear=True)

        # Threat section
        _, tts = add_textbox(sl, tx + 0.1, 1.9, 2.9, 0.25)
        para(tts, "THREAT:", size=12, bold=True, color=col, clear=True)
        _, ttb = add_textbox(sl, tx + 0.1, 2.15, 2.9, 1.5)
        ttb.word_wrap = True
        para(ttb, threat, size=13, color=C_WHITE, clear=True)

        # Arrow down
        _, tar = add_textbox(sl, tx + 1.2, 3.68, 0.7, 0.35)
        para(tar, "v", size=16, bold=True, color=col,
             align=PP_ALIGN.CENTER, clear=True)

        # Mitigation section
        add_divider(sl, tx + 0.1, 4.0, 2.9, color=col, height=0.03)
        _, tms = add_textbox(sl, tx + 0.1, 4.1, 2.9, 0.25)
        para(tms, "MITIGATION:", size=12, bold=True, color=col, clear=True)
        _, tmb = add_textbox(sl, tx + 0.1, 4.35, 2.9, 2.5)
        tmb.word_wrap = True
        para(tmb, mitigation, size=13, color=C_LIGHTGRAY, clear=True)

        tx += 3.25

    # Bottom note
    add_rect(sl, 0.35, 7.03, 12.63, 0.28, fill_color=C_BG2,
             line_color=C_GREEN, line_width_pt=0.8)
    _, tbot = add_textbox(sl, 0.5, 7.06, 12.3, 0.22)
    para(tbot, "Protocol duoc FREEZE truoc full run — giam researcher degrees of freedom",
         size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# S10 — TIMELINE & RESOURCES
# ═══════════════════════════════════════════════════════════════
def build_s10():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S10  Timeline & Resources",
                     "W7 pilot -> freeze gate -> W8 full -> W9 analysis",
                     speaker="Nghi", time_str="40s")
    slide_num_label(sl, 10)

    # Timeline
    weeks = [
        ("W5", "Freeze RQ,\nhypotheses,\nprompt", C_GRAY),
        ("W6", "GV approve\nDataset verify\nAPI setup", C_CYAN),
        ("W7\nPILOT", "5 Java + 5 Python\nFixture, logging\nRuntime check", C_AMBER),
        ("FREEZE\nGATE", "Manifest 50 units\nToolchain lock\nAmendment?", C_RED),
        ("W8\nFULL", "Full API run\nRandoop baseline\nMetrics all 50", C_GREEN),
        ("W9", "Pre-registered\nanalysis + CI\nDefense prep", C_PURPLE),
    ]
    wx = 0.35
    for label, desc, col in weeks:
        is_gate = "GATE" in label
        h = 1.6 if not is_gate else 1.9
        add_rect(sl, wx, 1.42, 1.98, h, fill_color=C_BG2,
                 line_color=col, line_width_pt=2.0 if is_gate else 1.2)
        add_rect(sl, wx, 1.42, 1.98, 0.42, fill_color=col)
        _, tw = add_textbox(sl, wx + 0.06, 1.44, 1.87, 0.36)
        para(tw, label, size=13, bold=True, color=C_BG,
             align=PP_ALIGN.CENTER, clear=True)
        _, td = add_textbox(sl, wx + 0.06, 1.92, 1.87, h - 0.55)
        td.word_wrap = True
        para(td, desc, size=12, color=C_LIGHTGRAY,
             align=PP_ALIGN.CENTER, clear=True)
        if wx < 11.5:
            _, ta = add_textbox(sl, wx + 2.0, 2.15, 0.25, 0.35)
            para(ta, ">", size=14, bold=True, color=C_GRAY,
                 align=PP_ALIGN.CENTER, clear=True)
        wx += 2.18

    # Resources + Cost card
    add_rect(sl, 0.35, 3.32, 6.0, 3.82, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=1.2)
    add_rect(sl, 0.35, 3.32, 6.0, 0.4, fill_color=C_CYAN)
    _, trh = add_textbox(sl, 0.5, 3.35, 5.7, 0.34)
    para(trh, "RESOURCES", size=14, bold=True, color=C_BG,
         align=PP_ALIGN.CENTER, clear=True)
    _, trc = add_textbox(sl, 0.5, 3.82, 5.7, 3.2)
    trc.word_wrap = True
    res = [
        ("API Budget", "GPT-4o snapshot gpt-4o-2024-08-06"),
        ("Cost estimate", "~$1.14 USD  (60 requests, ~87,500 tokens)"),
        ("Compute", "CPU + Git repo luu artifacts"),
        ("Input price", "$5.00 / 1M tokens"),
        ("Output price", "$15.00 / 1M tokens"),
        ("Buffer", "+20%  =>  tong ~$1.37 USD  (<$5 limit)"),
    ]
    first = True
    for k, v in res:
        if first:
            kv_line(trc, k, v, key_size=13, val_size=12, space=0)
            first = False
        else:
            kv_line(trc, k, v, key_size=13, val_size=12, space=6)

    # Contingency card
    add_rect(sl, 6.7, 3.32, 6.28, 3.82, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.2)
    add_rect(sl, 6.7, 3.32, 6.28, 0.4, fill_color=C_AMBER)
    _, tch = add_textbox(sl, 6.85, 3.35, 5.9, 0.34)
    para(tch, "CONTINGENCY PLAN", size=14, bold=True, color=C_BG,
         align=PP_ALIGN.CENTER, clear=True)
    _, tcc = add_textbox(sl, 6.85, 3.82, 5.9, 3.2)
    tcc.word_wrap = True
    conts = [
        "Model snapshot unavailable\n-> Dung, lap amendment, xin GV duyet",
        "Tool/fixture failure\n-> Fix + re-pilot, khong skip pilot",
        "API rate limit\n-> Chia batch, chay qua dem",
        "Pilot phat hien van de ky thuat\n-> Amendment trong 24h sau pilot meeting",
        "N thuc te < proposal\n-> Cap nhat §5.2 + power analysis",
    ]
    first = True
    for c in conts:
        if first:
            bullet(tcc, c, size=12, space_before=0)
            first = False
        else:
            bullet(tcc, c, size=12, space_before=5)


# ═══════════════════════════════════════════════════════════════
# S11 — CONTRIBUTION & Q&A
# ═══════════════════════════════════════════════════════════════
def build_s11():
    sl = new_slide()
    set_bg(sl, C_BG)
    add_slide_header(sl, "S11  Contribution & Q&A",
                     "Ket qua du kien va phong van GV",
                     speaker="Nghi", time_str="25s")
    slide_num_label(sl, 11)

    # Contributions column (left)
    add_rect(sl, 0.35, 1.42, 6.8, 5.6, fill_color=C_BG2,
             line_color=C_CYAN, line_width_pt=1.5)
    add_rect(sl, 0.35, 1.42, 6.8, 0.4, fill_color=C_CYAN)
    _, tch = add_textbox(sl, 0.5, 1.44, 6.5, 0.34)
    para(tch, "EXPECTED CONTRIBUTIONS", size=14, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)

    contribs = [
        ("Controlled protocol\nCC 5-15",
         "First study to evaluate GPT-4o on balanced Java/Python\nfiltered by CC 5-15 — bounded claim",
         C_CYAN),
        ("Paired Java comparison\nvs Randoop 4.3.4",
         "Paired mutation score comparison with statistical test,\neffect size, 95% CI & multiplicity correction",
         C_AMBER),
        ("Reproducibility package\n50/50 traceability",
         "Frozen manifest, prompt hash, raw logs, analysis script\n& amendment log — fully reproducible",
         C_GREEN),
    ]
    ry = 1.95
    for title, desc, col in contribs:
        add_rect(sl, 0.5, ry, 6.5, 1.48, fill_color=C_BG,
                 line_color=col, line_width_pt=1.0)
        add_rect(sl, 0.5, ry, 0.12, 1.48, fill_color=col)
        _, tt = add_textbox(sl, 0.72, ry + 0.05, 6.15, 0.5)
        para(tt, title, size=14, bold=True, color=col, clear=True)
        _, td = add_textbox(sl, 0.72, ry + 0.55, 6.15, 0.88)
        td.word_wrap = True
        para(td, desc, size=12, color=C_LIGHTGRAY, clear=True)
        ry += 1.6

    # NOT claiming column
    add_rect(sl, 7.5, 1.42, 5.48, 2.1, fill_color=RGBColor(0x1A, 0x10, 0x10),
             line_color=C_RED, line_width_pt=1.5)
    add_rect(sl, 7.5, 1.42, 5.48, 0.4, fill_color=C_RED)
    _, tnc = add_textbox(sl, 7.65, 1.44, 5.1, 0.34)
    para(tnc, "KHONG TUYEN BO", size=14, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)
    _, tncb = add_textbox(sl, 7.65, 1.92, 5.1, 1.55)
    tncb.word_wrap = True
    noclaims = [
        "GPT-4o tot hon trong moi truong hop",
        "CC 5-15 dai dien toan nganh",
        "Inferential stats la chua tung dung",
        "Ket qua pilot = ket qua final",
    ]
    for nc in noclaims:
        bullet(tncb, nc, size=13, color=C_RED, marker="x", space_before=3)

    # Q&A panel
    add_rect(sl, 7.5, 3.72, 5.48, 3.3, fill_color=C_BG2,
             line_color=C_AMBER, line_width_pt=1.5)
    add_rect(sl, 7.5, 3.72, 5.48, 0.4, fill_color=C_AMBER)
    _, tqa = add_textbox(sl, 7.65, 3.74, 5.1, 0.34)
    para(tqa, "Q&A — AI TRA LOI GI?", size=14, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER, clear=True)

    qas = [
        ("75% lay tu dau?", "Vinh", C_CYAN),
        ("Model, API, SDK, prompt?", "Thai", C_GREEN),
        ("RQ, H0/H1, stats, Randoop?", "Quan", C_AMBER),
        ("Python toolchain, threats?", "Van", C_PURPLE),
        ("Timeline, resources?", "Nghi", C_RED),
    ]
    _, tqab = add_textbox(sl, 7.62, 4.22, 5.3, 2.7)
    tqab.word_wrap = True
    first = True
    for q, who, col in qas:
        p = tqab.paragraphs[0] if first else tqab.add_paragraph()
        if first:
            p.clear()
            first = False
        p.space_before = Pt(4)
        r1 = p.add_run()
        r1.text = f"{q}  -> "
        r1.font.size = Pt(12)
        r1.font.color.rgb = C_LIGHTGRAY
        r1.font.name = "Calibri"
        r2 = p.add_run()
        r2.text = who
        r2.font.size = Pt(12)
        r2.font.bold = True
        r2.font.color.rgb = col
        r2.font.name = "Calibri"

    # Center Q&A invitation
    _, tqainv = add_textbox(sl, 2.5, 7.0, 8.33, 0.38)
    para(tqainv, "Nhom em xin nhan cau hoi cua Thay/Co",
         size=20, bold=True, color=C_CYAN, align=PP_ALIGN.CENTER, clear=True)


# ═══════════════════════════════════════════════════════════════
# BUILD ALL SLIDES
# ═══════════════════════════════════════════════════════════════
print("Building slides...")
build_s1()
print("  S1 Title & Scope done")
build_s2()
print("  S2 Problem -> GAP done")
build_s3()
print("  S3 Key Evidence done")
build_s4()
print("  S4 RQ & Statistical Tests done")
build_s5()
print("  S5 Pipeline & Model Spec done")
build_s6()
print("  S6 Dataset & Baseline done")
build_s7()
print("  S7 Work Packages done")
build_s8()
print("  S8 Evaluation Plan done")
build_s9()
print("  S9 Threats to Validity done")
build_s10()
print("  S10 Timeline & Resources done")
build_s11()
print("  S11 Contribution & Q&A done")

prs.save(str(OUT))
print(f"\n[DONE] Saved: {OUT}")
