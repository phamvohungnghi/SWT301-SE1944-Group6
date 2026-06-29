# -*- coding: utf-8 -*-
"""
update_slide2_add_numbers.py
----------------------------
Cập nhật Slide 2 (Problem → GAP) trong file PPT:
    GPT-4o_RBL3_Proposal_Defense_Visual_Revised.pptx

Nội dung mới theo existing-slide-revision-plan.md — S2:
  - Prior work (con số từ evidence table)
  - Confounders (Dataset khác nhau / Model • Prompt • Metric • Repair)
  - Bounded GAP (chưa có mẫu cân bằng Java/Python lọc trước CC 5–15)
  - Controlled protocol callout

Script tìm textbox lớn nhất trên Slide 2 để thay nội dung,
giữ nguyên các shape khác (nền, hình ảnh, …).

Yêu cầu: python-pptx  (pip install python-pptx)
"""

import pathlib
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Đường dẫn ──────────────────────────────────────────────
#   File gốc (KHÔNG dùng ~$... vì đó là lock file tạm của PowerPoint)
ppt_in = pathlib.Path(
    r"D:\pal\RBL-SWT-main_3\SWT301-SE1944-Group6\presentation"
    r"\GPT-4o_RBL3_Proposal_Defense_Visual_Revised.pptx"
)
ppt_out = pathlib.Path(
    r"D:\pal\RBL-SWT-main_3\SWT301-SE1944-Group6\presentation"
    r"\GPT-4o_RBL3_Proposal_Defense_Visual_Final.pptx"
)

if not ppt_in.exists():
    raise FileNotFoundError(f"Không tìm thấy file PPT gốc:\n  {ppt_in}")

# ── Mở presentation ────────────────────────────────────────
prs = Presentation(str(ppt_in))

# Slide 2 (index 1)
slide = prs.slides[1]


# ── Helper: thêm run với format riêng ──────────────────────
def add_run(paragraph, text, *, font_name="Calibri", size=20,
            bold=False, italic=False, color=(255, 255, 255)):
    """Thêm một run vào paragraph với format cụ thể."""
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor(*color)
    return run


def add_heading(tf, text, *, size=24, color=(0, 188, 212)):
    """Thêm dòng tiêu đề (Calibri Bold, cyan mặc định)."""
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    add_run(p, text, size=size, bold=True, color=color)
    return p


def add_bullet(tf, main_text, source_text=None, *, level=1):
    """
    Thêm bullet với phần nội dung chính (trắng) và phần nguồn (italic, xám nhạt).
    """
    p = tf.add_paragraph()
    p.level = level
    p.alignment = PP_ALIGN.LEFT
    # Bullet character
    add_run(p, "• ", size=20, color=(255, 255, 255))
    # Nội dung chính
    add_run(p, main_text, size=20, color=(255, 255, 255))
    # Nguồn (nếu có)
    if source_text:
        add_run(p, f"  {source_text}", size=18, italic=True,
                color=(204, 204, 204))
    return p


# ── Tìm textbox lớn nhất trên Slide 2 ─────────────────────
#    Để chỉ thay nội dung textbox chính, không ghi đè các shape khác.
target_shape = None
max_area = 0

for shape in slide.shapes:
    if shape.has_text_frame:
        area = shape.width * shape.height
        if area > max_area:
            max_area = area
            target_shape = shape

if target_shape is None:
    raise RuntimeError("Không tìm thấy textbox nào trên Slide 2!")

# ── Xoá nội dung cũ và ghi nội dung mới ───────────────────
tf = target_shape.text_frame
tf.clear()
tf.word_wrap = True

# --- Phần 1: Prior work ─────────────────────────────────────
add_heading(tf, "Prior work")
add_bullet(tf, "TestPilot: median branch coverage 52.8 %",
           "(Table 2, p. 10)")
add_bullet(tf, "AgoneTest: highest branch coverage 79.8 %",
           "(Fig 4, p. 12)")

# --- Phần 2: Confounders ────────────────────────────────────
p_spacer = tf.add_paragraph()  # khoảng cách
add_run(p_spacer, " ", size=10)

add_heading(tf, "Confounders")
add_bullet(tf, "Dataset khác nhau")
add_bullet(tf, "Model • Prompt • Metric • Repair")

# --- Phần 3: Bounded GAP ────────────────────────────────────
p_spacer = tf.add_paragraph()
add_run(p_spacer, " ", size=10)

add_heading(tf, "Bounded GAP", color=(255, 193, 7))  # amber/vàng nổi bật
add_bullet(tf, "Trong 10 papers: chưa có mẫu cân bằng Java/Python "
              "lọc trước CC 5–15")

# --- Callout: Controlled protocol ────────────────────────────
p_spacer = tf.add_paragraph()
add_run(p_spacer, " ", size=10)

p_callout = tf.add_paragraph()
p_callout.alignment = PP_ALIGN.CENTER
add_run(p_callout, "→ Controlled protocol", size=22, bold=True,
        color=(0, 230, 118))  # green accent

# ── Lưu file ──────────────────────────────────────────────
prs.save(str(ppt_out))
print(f"[OK] Slide 2 da duoc cap nhat va luu tai:\n   {ppt_out}")
